from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- CONFIGURATION ---
# Hyundai Branding Colors (for alignment)
HYUNDAI_BLUE = RGBColor(0, 44, 95)     
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(100, 100, 100)
# Invincible Ocean Text Color
IO_TEXT = RGBColor(40, 40, 40) 

def create_full_presentation():
    print("Generating Full Presentation...")
    prs = Presentation()
    
    # 1. HELPER: Title Slide
    def add_title_slide(title_text, subtitle_text, author_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = HYUNDAI_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        title_tf = title_box.text_frame
        title_tf.text = title_text
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(54)
        title_p.font.color.rgb = WHITE
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.text = subtitle_text
        sub_p = sub_tf.paragraphs[0]
        sub_p.font.size = Pt(24)
        sub_p.font.color.rgb = RGBColor(200, 200, 200)
        sub_p.alignment = PP_ALIGN.CENTER

        # Author
        auth_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        auth_tf = auth_box.text_frame
        auth_tf.text = author_text
        auth_p = auth_tf.paragraphs[0]
        auth_p.font.size = Pt(18)
        auth_p.font.color.rgb = RGBColor(180, 210, 255)
        auth_p.alignment = PP_ALIGN.CENTER

    # 2. HELPER: Text Content Slide
    def add_content_slide(title_text, content_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = HYUNDAI_BLUE
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = IO_TEXT
            p.space_after = Pt(14)

    # 3. HELPER: Diagram Slide (Split Layout)
    def add_diagram_slide(title_text, content_lines, placeholder_label):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = HYUNDAI_BLUE
        title.text_frame.paragraphs[0].font.bold = True
        
        # Left Side Text
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.space_after = Pt(12)
            
        # Right Side Placeholder (Grey Box)
        shape = slide.shapes.add_shape(
            1, Inches(5.2), Inches(1.5), Inches(4.5), Inches(4.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        
        p = shape.text_frame.paragraphs[0]
        p.text = f"[PASTE {placeholder_label} HERE]"
        p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.CENTER

    # 4. HELPER: Closing Slide
    def add_closing_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = HYUNDAI_BLUE
        
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(60)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        c_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        c_tf = c_box.text_frame
        c_p = c_tf.paragraphs[0]
        c_p.text = "Invincible Ocean"
        c_p.font.size = Pt(24)
        c_p.font.color.rgb = RGBColor(180, 210, 255)
        c_p.alignment = PP_ALIGN.CENTER

    # --- SLIDE CONTENT ---

    # Slide 1
    add_title_slide("AI SALES BUDDY", "Technical Solution Proposal", "By Invincible Ocean")

    # Slide 2
    add_content_slide("1. Solution Overview", [
        "The AI Sales Buddy is an intelligent assistant embedded in the LMS and H-Smart App.",
        "• Empower Sales Consultants with instant product knowledge.",
        "• Provide real-time competitive intelligence.",
        "• Support 6 Indian languages natively.",
        "• Deliver coaching and objection handling on the fly."
    ])

    # Slide 3
    add_diagram_slide("2. System Architecture", [
        "Modular, Enterprise-Grade Design:",
        "• Presentation Layer: LMS & H-Smart Widgets.",
        "• Gateway Layer: SSO Auth & Rate Limiting.",
        "• Intelligence Layer: Orchestrator & Agents.",
        "• Connector Layer: LMS API, Web Search.",
        "• Data Layer: Vector Store & Cache."
    ], "ARCHITECTURE DIAGRAM")

    # Slide 4
    add_diagram_slide("3. Data Flow Strategy", [
        "How a query is processed:",
        "1. Query Understanding: Detect intent (e.g., 'Compare').",
        "2. Parallel Retrieval: Fetch internal Hyundai data + External Web data.",
        "3. Data Merger: Normalize prices and specs.",
        "4. Intelligent Response: Highlight Hyundai USPs.",
    ], "DATA FLOW DIAGRAM")

    # Slide 5
    add_content_slide("4. Competitive Intelligence", [
        "Real-time advantage over static PDFs:",
        "• Data Source: Live web search + cached competitor specs.",
        "• Comparison Engine: Automatically normalizes features.",
        "• USP Highlighting: Identifies where Hyundai wins.",
        "• Output: Generates clear comparison tables."
    ])

    # Slide 6
    add_diagram_slide("5. Vernacular Support", [
        "Native support for 6 Languages:",
        "• Hindi, Tamil, Telugu, Kannada, Malayalam, English.",
        "• Process: Speech-to-Text -> Translate -> Process -> Text-to-Speech.",
        "• Result: Consultants speak naturally in their local dialect."
    ], "VERNACULAR FLOW")

    # Slide 7
    add_content_slide("6. Security Architecture", [
        "Defense-in-depth approach:",
        "• Layer 1 (Edge): DDoS protection & Rate limiting.",
        "• Layer 2 (Auth): SSO Token validation & Session mgmt.",
        "• Layer 3 (Data): AES-256 Encryption & India data residency.",
        "• Layer 4 (Audit): Full logging of all queries for compliance."
    ])

    # Slide 8
    add_content_slide("7. Analytics & Insights", [
        "Dashboard for HO and Regional Offices:",
        "• Track top queries (e.g., 'What is ADAS?').",
        "• Identify training gaps in specific regions.",
        "• Monitor usage adoption per dealership.",
        "• Exportable reports for strategy planning."
    ])

    # Slide 9
    add_content_slide("8. Why Invincible Ocean?", [
        "• Production-Ready: Built on our proven automotive AI infrastructure.",
        "• Multi-Agent Expertise: Specialized agents for distinct sales tasks.",
        "• Non-Invasive: Simple iFrame integration (No major IT changes).",
        "• Speed: 15-week delivery timeline.",
        "• Scalability: Designed for 10,000+ concurrent users."
    ])

    # Slide 10
    add_closing_slide()

    prs.save('Invincible_Ocean_Proposal_Full.pptx')
    print("SUCCESS: Full presentation 'Invincible_Ocean_Proposal_Full.pptx' created!")

if __name__ == "__main__":
    create_full_presentation()